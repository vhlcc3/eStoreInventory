# -*- coding: utf-8 -*-
"""
Created on Wed Jun 15 11:31:55 2022

@author: volke

Thoughts -

    * Alignment of captions can be improved
    * Allow content of captions to be via a function that can be overloaded
    * For multiple columns of items, wbn (would be nice) if left column were
      left justified, right column right justified and others center justified.


"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
from PIL.ExifTags import TAGS
import os
import pandas as pd

from datetime import datetime

class pptxCatalog():
    
    def __init__( self, pptxPath, pptxConfig ):
            
        self.imgRotations = { 0: 0, 1: 0, 3: 180, 6: 270, 8: 90 }

        self.iC = Presentation()
        self.cfg = {
                    "SlideWidth"   : 8.5,
                    "SlideHeight"  : 11.0,
                    "CaptionHt" : 0.4,
                    "ImageOffsetLeft"     : 0.5,
                    "ImageOffsetTop"      : 0.8,
                    "ImageGap"      : 0.4,
                    "CaptionGap"    : 0.0,
                    "SlidesAcross"   : 2,
                    "SlidesDown"     : 3,
                    "CaptionFontSize"     : 8,
                    "ErrorFontSize"       : 6,
                    "CaptionBold"         : True,
                    "ErrorBold"           : False

                    }
        # Note that the 'Inches()' method converts actual float inches to the integer
        # inch representation required by pypptx, which is Inches * 1,000,000
        
        self.rescale = 3
        if not pptxConfig is None:
            for cfgItem in pptxConfig:
                self.cfg[cfgItem] = pptxConfig[cfgItem]
        self.slWidth= self.cfg["SlideWidth"]
        self.slHeight= self.cfg["SlideHeight"]
        self.iC.slide_width = Inches( self.slWidth )
        self.iC.slide_height= Inches( self.slHeight )
        self.iLeft  = self.cfg["ImageOffsetLeft"]
        self.left   = self.iLeft
        self.iTop   = self.cfg["ImageOffsetTop"]
        self.top    = self.iTop
        self.iGap   = self.cfg["ImageGap"]
        self.sAcross= self.cfg["SlidesAcross"]
        self.sDown  = self.cfg["SlidesDown"]
        self.captionHt = self.cfg["CaptionHt"]
        self.errorFrame = self._addErrorSlide()
        self.pageNo = 0
        self.newPage= True
        self.iWidth = (self.slWidth - ((self.iGap * (self.sAcross-1)) + 2.0*self.iLeft)) / float(self.sAcross)
        self.iHeight= (self.slHeight - ((self.iGap*(self.sDown-1)) + 2.0*self.iTop)) / float(self.sDown)
        self.inchImgWid = Inches( self.iWidth )
        self.inchImgHt  = Inches( self.iHeight )
        self.iRatio = self.iHeight / self.iWidth
        self.pptxPath = pptxPath
    
    def imageSize( self ):
        return ( self.iWidth, self.iHeight )
        
    def addItemWithImage( self, itemInfo, imageId, imageFile, imageRatio ):

        if self.newPage:
            self.pageNo += 1
            print( "Page %d" %(self.pageNo))
            self.slide = self.iC.slides.add_slide(self.iC.slide_layouts[6])
            self.newPage = False

        try:
            print("Image %s" % (imageFile))
            with open( imageFile, mode = "br" ) as iFs:
                print(" Adding %s" %(imageFile))
                if imageRatio > self.iRatio:
                    imW = None
                    imH = self.inchImgHt
                    tbVertical = Inches( self.top + self.iHeight + self.cfg["CaptionGap"])
                else:
                    imW = self.inchImgWid
                    imH = None
                    tbVertical = Inches( self.top + ( self.iWidth * imageRatio ) + self.cfg["CaptionGap"])
                self.slide.shapes.add_picture( iFs, Inches(self.left), Inches(self.top) , width=imW, height=imH )
                sTextBox = self.slide.shapes.add_textbox( Inches(self.left), tbVertical, self.inchImgWid, Inches(self.captionHt))
                pceName  = itemInfo["Name"]
                self._imgCaption( sTextBox, 
                           pceName,
                           str(itemInfo["Subtype"]),
                           str(itemInfo["Catalogue Price"]),
                           imageId + " / " + str(itemInfo["ID"])  )
                self.left += self.iWidth + self.iGap
                if (self.left+self.iWidth) > self.slWidth:
                    self.left = self.iLeft
                    self.top += self.iHeight + self.iGap
                    if (self.top+self.iHeight) > self.slHeight:
                        self.top = self.iTop
                        self.newPage = True

        except OSError as error:
            print("Error accessing %s:\n%s" %(imageFile,error))
            self._imgError( self.errorFrame, str(error) )

    # Create a text frame for capturing error messages 
    
    def _addErrorSlide( self ):
        eS = self.iC.slides.add_slide(self.iC.slide_layouts[6])
        eB = eS.shapes.add_textbox( Inches(0.5), Inches(0.5 ), Inches(7.5), Inches(9.0))
        self._imgError( eB.text_frame, "Errors and Warnings - %s" %(str(datetime.now())))
        return eB.text_frame

    def _imgError( self, tF, errorMsg ):
        tFp = tF.add_paragraph()
        run = tFp.add_run()
        run.text = errorMsg

    # Add text to a text box in a given font
    
    def _imgCaption( self, tB, iName, iType, iPrice, iIndexes ):
        tF = tB.text_frame
        # Textbox itself seems to be aligned on its center so am aligning contained
        # text in same way - must be some way of positioning text box to the left ...
        # Maybe by padding the text with trailing blanks to the required length.
        tF.paragraphs[0].alignment = PP_ALIGN.CENTER
        p  = tF.paragraphs[0]
        run = p.add_run()
        run.text = iName
        if iType and iType.strip():
            run.text = iName + " (" + iType + ")" + " " + iPrice
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(self.cfg["CaptionFontSize"])
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        run = p.add_run()
        run.text = "\n" + iIndexes
        
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(self.cfg["ErrorFontSize"])
        font.bold = self.cfg["ErrorBold"]
        font.italic = None  # cause value to be inherited from theme
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
        

    def saveOutput( self ):
        self.iC.save( self.pptxPath )
        # self.iC.save( "outputs" + os.path.sep + "Image_Cat.pptx")       


    